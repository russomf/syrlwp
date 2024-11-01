# report_pptx.py
# Streamlining Your Research Laboratory with Python
# Authors:   Mark F. Russo, Ph.D and William Neil 
# Publisher: John Wiley & Sons, Inc.
# License:   MIT (https://opensource.org/licenses/MIT)

from io   import BytesIO
from pptx import Presentation
from pptx.util import Inches

# Save a matplotlib Figure to a BytesIO and return
def fig2BytesIO(fig):
    bfile = BytesIO()
    fig.savefig(bfile, format='png')
    bfile.seek(0)
    return bfile

def generate_pptx(figs, lol, fname):
    # Create new Presentation
    prs = Presentation()
    
    # Get useful layouts and create Length objects
    title_layout   = prs.slide_layouts.get_by_name("Title Slide")
    content_layout = prs.slide_layouts.get_by_name("Title and Content")
    in1 = Inches(1)
    in2 = Inches(2)
    in3 = Inches(3)
    in4 = Inches(4)
    in6 = Inches(6)
    
    # Add Title slide and title
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Experimental Data"
    
    # Generate all figures and add each to a separate slide
    for fig in figs:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.add_picture(fig2BytesIO(fig), in2, in1, width=in6)

    # Add final slide with generated table
    slide = prs.slides.add_slide(content_layout)
    nrows = len(lol)                        # Get Table size and add Table
    ncols = max([len(row) for row in lol])
    frame = slide.shapes.add_table(nrows, ncols, in3, in2, in4, in4)
    tbl   = frame.table                     # Get a Table reference
    
    for r in range(nrows):                  # Add data to table cells
        for c in range(ncols):
            tbl.cell(r, c).text = str(lol[r][c])

    # Generate presentation and save
    prs.save(fname)

if __name__ == '__main__':
    from report_utils import *

    # Read data and transpose so columns hold data for single samples
    df = pd.read_csv("data_plate.csv", header=None)
    df = df.transpose()
    
    # Generate heatmap Figure from data
    fig1 = plot_heatmap(df, "Experimental Data")
    figs = [fig1]

    # Fit all data to four-parameter logistic model
    guess  = [0, 2, -5, 1]                            # Parameter guesses
    bounds = [[0,0,-9,0], [1,1000,0,1]]               # Parameter bounds
    concs  = [-9, -8, -7, -6, -5, -4, -3, -2, -1, 0]  # log(concentrations)
    params = fit_4PLs(df, concs, guess, bounds)       # Fit all data

    # Create eight Figures with a single plot for data and parameter set.
    for col in range(8):
        title = f"Sample {col+1}"
        fig = plot_fit(df[col], concs, title, params[col])
        figs.append(fig)
    
    # Collect IC50 results as list of lists
    IC50s = [f'{pow(10, float(pset[2])):0.3E}' for pset in params]
    lol   = [ (f'Sample {i+1}', IC50s[i]) for i in range(8) ]
    lol.insert(0, ["Sample ID", "IC50 [M]"])          # Add column headers
    
    # Generate presentation document
    generate_pptx(figs, lol, 'report.pptx')
